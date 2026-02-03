import os
import torch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration


# alt text generation is not very good at the moment...

device = "cuda" if torch.cuda.is_available() else "cpu"
processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
model = BlipForConditionalGeneration.from_pretrained(
    "Salesforce/blip-image-captioning-base"
).to(device)


def generate_alt_text(image_path):
    image = Image.open(image_path).convert("RGB")
    inputs = processor(image, return_tensors="pt").to(device)
    output = model.generate(**inputs, max_new_tokens=50)
    return processor.decode(output[0], skip_special_tokens=True)


def accessibility_processor(directory):

    for filename in os.listdir(directory):

        if filename.endswith(".pptx"):

            filepath = os.path.join(directory, filename)
            prs = Presentation(filepath)

            filepath_alt_text = os.path.join(new_directory, filename)
            alt_text_output_file = f"{os.path.splitext(filepath_alt_text)[0]}_alt_text"

            with open(alt_text_output_file, 'w', encoding='utf-8') as f:

                for slide_idx, slide in enumerate(prs.slides):

                    for shape_idx, shape in enumerate(slide.shapes):

                        if "Title" in shape.name:
                            cursor_sp = slide.shapes[0]._element
                            cursor_sp.addprevious(shape._element)

                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                            image_ext = "png"
                            image_filename = f"_tmp_slide{slide_idx}_shape{shape_idx}.{image_ext}"
                            image_path = os.path.join(new_directory, image_filename)

                            with open(image_path, "wb") as img_file:
                                img_file.write(shape.image.blob)

                            alt_text = generate_alt_text(image_path)

                            f.write(
                                f"Slide {slide_idx + 1}, Image {shape_idx + 1}: {alt_text}\n"
                            )

                            os.remove(image_path)

            new_filepath = os.path.join(
                new_directory,
                f"{os.path.splitext(filename)[0]}_updated.pptx"
            )

            prs.save(new_filepath)

            new_filepath_original_file = os.path.join(new_directory, filename)
            os.rename(filepath, new_filepath_original_file)


directory_path = "C:/Users/nairt/Desktop/GitHub Repos/UW_Accessibility_Enhancer"
new_directory = "C:/Users/nairt/Desktop/GitHub Repos/UW_Accessibility_Enhancer"

accessibility_processor(directory_path)
